# src/core/ppt_tools/image_handler.py
import io
import logging
from pathlib import Path
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw

logger = logging.getLogger(__name__)

_EMU_PER_PX = 9525  # 96 DPI: 914400 EMU/inch ÷ 96 px/inch


def _make_placeholder(width_emu: int, height_emu: int, missing_path: Path | str) -> io.BytesIO:
    """Generate an in-memory placeholder PNG for a missing image."""
    px_w = max(60, round(width_emu / _EMU_PER_PX))
    px_h = max(40, round(height_emu / _EMU_PER_PX))

    img = Image.new("RGB", (px_w, px_h), color=(255, 220, 220))
    draw = ImageDraw.Draw(img)

    # Border
    draw.rectangle([0, 0, px_w - 1, px_h - 1], outline=(180, 0, 0), width=3)

    # Text: label + filename (truncated to fit)
    filename = Path(missing_path).name
    label = "Image not found:"
    max_chars = max(10, px_w // 7)
    if len(filename) > max_chars:
        filename = filename[:max_chars - 1] + "…"

    try:
        from PIL import ImageFont
        font = ImageFont.load_default(size=max(10, px_h // 8))
    except TypeError:
        from PIL import ImageFont
        font = ImageFont.load_default()

    draw.text((6, 6), label, fill=(140, 0, 0), font=font)
    draw.text((6, 6 + px_h // 6), filename, fill=(140, 0, 0), font=font)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

def _find_shape_recursive(shapes, shape_name: str):
    """
    ฟังก์ชันภายใน: ค้นหา Shape ตามชื่อแบบเจาะลึกเข้าไปใน Group
    """
    for shape in shapes:
        if shape.name == shape_name:
            return shape
        # ถ้าเจอ Group, ให้เจาะเข้าไปหารอบใหม่
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found_shape = _find_shape_recursive(shape.shapes, shape_name)
            if found_shape:
                return found_shape
    return None

def replace_image_by_name(slide: Slide, shape_name: str, new_image_path: Path | str) -> bool:
    """
    ค้นหาและแทนที่รูปภาพตามชื่อ โดยคงตำแหน่ง ขนาด และสัดส่วนเดิมไว้
    * รองรับรูปภาพเดี่ยว และรูปภาพที่ถูกจัดกลุ่ม (Group Shape) *
    * หมายเหตุ: รูปใหม่จะถูกนำมาวางไว้ชั้นบนสุด (Layering) *
    """
    # 1. ตรวจสอบไฟล์รูปภาพใหม่
    if not new_image_path:
        logger.warning(f"ไม่ได้ระบุ Path สำหรับรูปภาพใหม่ ของ Shape: '{shape_name}'")
        return False

    path_obj = Path(new_image_path)
    use_placeholder = not path_obj.exists()
    if use_placeholder:
        logger.error(f"ไม่พบไฟล์รูปภาพใหม่ตาม Path: {new_image_path} (สำหรับ Shape: '{shape_name}') — แทรก placeholder แทน")

    # 2. ค้นหา Shape เป้าหมาย (แบบเจาะลึก Group)
    target_shape = _find_shape_recursive(slide.shapes, shape_name)

    # 3. ตรวจสอบ Error
    if not target_shape:
        logger.warning(f"ไม่พบ Object ชื่อ '{shape_name}' บนสไลด์หน้า {slide.slide_id}")
        return False

    if target_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        logger.warning(f"Object ชื่อ '{shape_name}' ถูกพบ แต่ไม่ใช่รูปภาพ (เป็นชนิด: {target_shape.shape_type})")
        return False

    if not use_placeholder:
        logger.info(f"กำลังแทนที่รูปภาพ: '{shape_name}' ด้วยไฟล์: {path_obj.name}")

    # 4. เก็บค่าตำแหน่งและขนาดเดิม
    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    # 5. ลบรูปเดิมออก
    old_picture_element = target_shape.element
    old_picture_element.getparent().remove(old_picture_element)

    # 6. แทรกรูปใหม่ (หรือ placeholder) ลงไปที่ตำแหน่งและขนาดเดิม
    try:
        image_source = _make_placeholder(width, height, path_obj) if use_placeholder else str(path_obj)
        pic = slide.shapes.add_picture(image_source, left, top, width, height)
    except Exception as e:
        logger.error(f"เกิดข้อผิดพลาดขณะแทรกรูปภาพใหม่ '{shape_name}': {e}")
        return False

    # 7. ส่งรูปไปด้านหลังสุด (Send to Back)
    sp = pic._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)

    return not use_placeholder