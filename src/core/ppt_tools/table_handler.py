# src/core/ppt_tools/table_handler.py
import logging
import math
from copy import deepcopy
from lxml import etree
from pptx.slide import Slide

logger = logging.getLogger(__name__)

_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A   = f"{{{_NS}}}"
_FONT = "DB Heavent"
_FONT_PANOSE = "02000506060000020004"

# Run-level attributes matching the template
_RPR_ATTRS = (
    'kumimoji="0" sz="1050" b="0" i="0" u="none" '
    'strike="noStrike" kern="1200" cap="none" spc="0" '
    'normalizeH="0" baseline="0" noProof="0"'
)
_RPR_FONT_ELEMS = (
    f'<a:ln><a:noFill/></a:ln>'
    f'<a:effectLst/><a:uLnTx/><a:uFillTx/>'
    f'<a:latin typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
    f'<a:ea typeface="+mn-ea"/>'
    f'<a:cs typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
)

# Cell styling
_FILL_POSITIVE  = f'<a:solidFill xmlns:a="{_NS}"><a:srgbClr val="31A1C0"><a:alpha val="20000"/></a:srgbClr></a:solidFill>'
_FILL_NEGATIVE  = f'<a:solidFill xmlns:a="{_NS}"><a:srgbClr val="F5EAAB"/></a:solidFill>'
_FILL_NEAR_ZERO = f'<a:solidFill xmlns:a="{_NS}"><a:schemeClr val="bg1"><a:lumMod val="95000"/></a:schemeClr></a:solidFill>'

_COLOR_POSITIVE  = '<a:srgbClr val="1D7381"/>'
_COLOR_NEGATIVE  = '<a:srgbClr val="9C551F"/>'
_COLOR_NEAR_ZERO = '<a:prstClr val="black"/>'

# ──────────────────────────────────────────────────────────
# Table layout normalization specs (all values in EMU)
# 1 inch = 914400 EMU
# ──────────────────────────────────────────────────────────

# Region tables (tbl_region_*): pages 8, 10, 12, 14
# Col widths: 0.95" (name) + 6 × 0.58" (months)
# Data row height: 0.31"  |  Header row (h=0): leave untouched
_REGION_COL_WIDTHS  = [868680, 530352, 530352, 530352, 530352, 530352, 530352]
_REGION_ROW_HEIGHT  = 283464   # 0.31"

# Basin tables (tbl_basin_*): pages 32, 34, 36, 38
# Col widths: 1.05" (name) + 6 × 0.59" (months)
# Header row height: 0.30"  |  Data row height: 0.36"
_BASIN_COL_WIDTHS        = [960120, 539496, 539496, 539496, 539496, 539496, 539496]
_BASIN_HEADER_ROW_HEIGHT = 274320   # 0.30"
_BASIN_DATA_ROW_HEIGHT   = 329184   # 0.36"


# ──────────────────────────────────────────────────────────
# XML element builders
# ──────────────────────────────────────────────────────────

def _xml_run(lang: str, text: str, color_inner: str) -> etree._Element:
    safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return etree.fromstring(
        f'<a:r xmlns:a="{_NS}">'
        f'<a:rPr lang="{lang}" {_RPR_ATTRS}>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:solidFill>{color_inner}</a:solidFill>'
        f'<a:effectLst/><a:uLnTx/><a:uFillTx/>'
        f'<a:latin typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
        f'<a:ea typeface="+mn-ea"/>'
        f'<a:cs typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
        f'</a:rPr>'
        f'<a:t>{safe}</a:t>'
        f'</a:r>'
    )


def _xml_br(color_inner: str) -> etree._Element:
    return etree.fromstring(
        f'<a:br xmlns:a="{_NS}">'
        f'<a:rPr lang="th-TH" {_RPR_ATTRS}>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:solidFill>{color_inner}</a:solidFill>'
        f'<a:effectLst/><a:uLnTx/><a:uFillTx/>'
        f'<a:latin typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
        f'<a:ea typeface="+mn-ea"/>'
        f'<a:cs typeface="{_FONT}" panose="{_FONT_PANOSE}" pitchFamily="2" charset="-34"/>'
        f'</a:rPr>'
        f'</a:br>'
    )


# ──────────────────────────────────────────────────────────
# Cell content builder
# ──────────────────────────────────────────────────────────

def _build_data_cell(tc: etree._Element, anomaly: float, pct: float) -> None:
    """
    Rewrite a single <a:tc> element's text and fill based on anomaly/pct values.
    Preserves existing <a:pPr> paragraph properties.
    """
    if math.isnan(anomaly) or math.isnan(pct):
        logger.warning("NaN value encountered in cell data — skipping cell.")
        return

    # Determine case
    if pct >= 1.0:
        color_inner = _COLOR_POSITIVE
        fill_xml    = _FILL_POSITIVE
    elif pct <= -1.0:
        color_inner = _COLOR_NEGATIVE
        fill_xml    = _FILL_NEGATIVE
    else:
        color_inner = _COLOR_NEAR_ZERO
        fill_xml    = _FILL_NEAR_ZERO

    # ── Update cell background fill in <a:tcPr> ──────────────
    tcPr = tc.find(f"{_A}tcPr")
    if tcPr is not None:
        for old_fill in tcPr.findall(f"{_A}solidFill"):
            tcPr.remove(old_fill)
        tcPr.append(etree.fromstring(fill_xml))

    # ── Rebuild paragraph content in <a:txBody><a:p> ─────────
    txBody = tc.find(f"{_A}txBody")
    if txBody is None:
        return

    p = txBody.find(f"{_A}p")
    if p is None:
        return

    # Preserve existing <a:pPr> (paragraph formatting)
    pPr = p.find(f"{_A}pPr")
    pPr_copy = deepcopy(pPr) if pPr is not None else None

    # Clear all children from <a:p>
    for child in list(p):
        p.remove(child)

    # Restore <a:pPr>
    if pPr_copy is not None:
        p.append(pPr_copy)

    # ── Build text content ────────────────────────────────────
    if -1.0 < pct < 1.0:
        # Single line: "ไม่เกิน ±1" (th-TH) + "%" (en-US)
        p.append(_xml_run("th-TH", "ไม่เกิน \u00b11", color_inner))
        p.append(_xml_run("en-US", "%",               color_inner))
    else:
        pct_str  = f"{int(round(pct)):+d}"
        # Line 1: mm value — use "ไม่เกิน ±1 มม." if anomaly rounds to zero
        if int(round(anomaly)) == 0:
            p.append(_xml_run("th-TH", "ไม่เกิน \u00b11 มม.", color_inner))
        else:
            p.append(_xml_run("en-US", f"{int(round(anomaly)):+d}", color_inner))
        p.append(_xml_br(color_inner))
        p.append(_xml_run("en-US", f"({pct_str}%)", color_inner))


# ──────────────────────────────────────────────────────────
# Table layout normalizer
# ──────────────────────────────────────────────────────────

def _normalize_table_layout(table, shape_name: str) -> None:
    """
    Enforce column widths and row heights based on shape name prefix.
    - tbl_region_*: 0.95" name col + 6×0.58" month cols; data rows 0.31"
    - tbl_basin_*:  1.05" name col + 6×0.59" month cols; header 0.30", data rows 0.36"
    Region header row (h=0) is intentionally skipped.
    """
    tbl = table._tbl

    if shape_name.startswith("tbl_region_"):
        col_widths = _REGION_COL_WIDTHS
        header_h   = None          # keep h=0 as-is
        data_h     = _REGION_ROW_HEIGHT

    elif shape_name.startswith("tbl_basin_"):
        col_widths = _BASIN_COL_WIDTHS
        header_h   = _BASIN_HEADER_ROW_HEIGHT
        data_h     = _BASIN_DATA_ROW_HEIGHT

    else:
        return  # unknown shape — no normalization

    # Set column widths
    for i, gc in enumerate(tbl.findall(f"{_A}gridCol")):
        if i < len(col_widths):
            gc.set("w", str(col_widths[i]))

    # Set row heights
    rows = tbl.findall(f"{_A}tr")
    for i, tr in enumerate(rows):
        if i == 0:
            if header_h is not None:
                tr.set("h", str(header_h))
            # else: leave untouched (region's h=0 stays)
        else:
            tr.set("h", str(data_h))


# ──────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────

def fill_rain_diff_table(
    slide: Slide,
    shape_name: str,
    name_to_data: dict,
) -> bool:
    """
    Fill a 2-column observed-vs-forecast diff table (Group 2.10.1–2.10.3).

    Args:
        slide:        Target slide.
        shape_name:   Name of the table shape (e.g. 'tbl_region_diff').
        name_to_data: {thai_name: {"anomaly": float, "percent": float}}

    Returns True if all data rows were matched and written successfully.
    """
    table_shape = next(
        (s for s in slide.shapes if s.name == shape_name and s.has_table),
        None,
    )
    if table_shape is None:
        logger.warning(f"Table shape '{shape_name}' not found on slide.")
        return False

    table = table_shape.table
    ok = True

    for row_idx in range(1, len(table.rows)):
        name_cell = table.cell(row_idx, 0)
        row_name  = name_cell.text_frame.text.strip()

        if row_name not in name_to_data:
            logger.warning(f"[{shape_name}] Name '{row_name}' not found in data — row skipped.")
            ok = False
            continue

        val = name_to_data[row_name]
        anomaly = float(val["anomaly"])
        pct     = float(val["percent"])

        tc = table.cell(row_idx, 1)._tc

        # Determine styling
        if pct >= 1.0:
            color_inner = _COLOR_POSITIVE
            fill_xml    = _FILL_POSITIVE
        elif pct <= -1.0:
            color_inner = _COLOR_NEGATIVE
            fill_xml    = _FILL_NEGATIVE
        else:
            color_inner = _COLOR_NEAR_ZERO
            fill_xml    = _FILL_NEAR_ZERO

        # Update cell background fill
        tcPr = tc.find(f"{_A}tcPr")
        if tcPr is not None:
            for old_fill in tcPr.findall(f"{_A}solidFill"):
                tcPr.remove(old_fill)
            tcPr.append(etree.fromstring(fill_xml))

        # Rebuild paragraph
        txBody = tc.find(f"{_A}txBody")
        if txBody is None:
            continue
        p = txBody.find(f"{_A}p")
        if p is None:
            continue

        pPr = p.find(f"{_A}pPr")
        pPr_copy = deepcopy(pPr) if pPr is not None else None
        for child in list(p):
            p.remove(child)
        if pPr_copy is not None:
            p.append(pPr_copy)

        if -1.0 < pct < 1.0:
            p.append(_xml_run("th-TH", "ไม่เกิน \u00b11", color_inner))
            p.append(_xml_run("en-US", "%",               color_inner))
        else:
            pct_str = f"{int(round(pct)):+d}"
            if int(round(anomaly)) == 0:
                p.append(_xml_run("th-TH", "ไม่เกิน \u00b11 มม.", color_inner))
            else:
                p.append(_xml_run("en-US", f"{int(round(anomaly)):+d}", color_inner))
            p.append(_xml_br(color_inner))
            p.append(_xml_run("en-US", f"({pct_str}%)", color_inner))

    return ok


def fill_rain_table(
    slide: Slide,
    shape_name: str,
    name_to_data: dict,
) -> bool:
    """
    Fill a PPT table shape with rainfall anomaly/percent data.

    Args:
        slide:        Target slide.
        shape_name:   Name of the table shape (e.g. 'tbl_basin_hii_left').
        name_to_data: {thai_name: [{"anomaly": float, "percent": float}, ...]}
                      List must have 6 entries (t1–t6, matching col indices 1–6).

    Returns True if all data rows were matched and written successfully.
    """
    table_shape = next(
        (s for s in slide.shapes if s.name == shape_name and s.has_table),
        None,
    )
    if table_shape is None:
        logger.warning(f"Table shape '{shape_name}' not found on slide.")
        return False

    table = table_shape.table
    _normalize_table_layout(table, shape_name)
    ok = True

    for row_idx in range(1, len(table.rows)):
        name_cell = table.cell(row_idx, 0)
        row_name  = name_cell.text_frame.text.strip()

        if row_name not in name_to_data:
            logger.warning(f"[{shape_name}] Name '{row_name}' not found in data — row skipped.")
            ok = False
            continue

        values = name_to_data[row_name]  # list of 6 dicts

        for col_idx, val in enumerate(values, start=1):
            if col_idx >= len(table.columns):
                break
            try:
                _build_data_cell(
                    tc      = table.cell(row_idx, col_idx)._tc,
                    anomaly = float(val["anomaly"]),
                    pct     = float(val["percent"]),
                )
            except Exception as e:
                logger.error(f"[{shape_name}] Cell [{row_idx + 1},{col_idx + 1}]: {e}")
                ok = False

    return ok
