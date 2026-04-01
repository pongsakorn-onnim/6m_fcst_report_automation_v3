# src/core/ppt_tools/text_handler.py
import logging
from copy import deepcopy
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

logger = logging.getLogger(__name__)

# ==========================================
# 1. Utilities จัดการวันที่ (จากโค้ดของคุณ)
# ==========================================
THAI_MONTHS = [
    "มกราคม", "กุมภาพันธ์", "มีนาคม", "เมษายน", "พฤษภาคม", "มิถุนายน",
    "กรกฎาคม", "สิงหาคม", "กันยายน", "ตุลาคม", "พฤศจิกายน", "ธันวาคม"
]

def get_thai_month(month_idx: int) -> str:
    """คืนค่าชื่อเดือนไทยเต็ม (1-12)"""
    if 1 <= month_idx <= 12:
        return THAI_MONTHS[month_idx - 1]
    return ""

def get_buddhist_year(year: int) -> int:
    """แปลง ค.ศ. เป็น พ.ศ."""
    return year + 543

def get_next_months(start_year: int, start_month: int, n: int):
    if n <= 0:
        return []
    results = []
    for i in range(n):
        future_month_idx = (start_month + i - 1) % 12 + 1
        year_offset = (start_month + i - 1) // 12
        future_year = start_year + year_offset
        results.append({
            "year": future_year,
            "month": future_month_idx,
            "thai_name": get_thai_month(future_month_idx),
            "buddhist_year": get_buddhist_year(future_year),
        })
    return results

def get_months_for_leads(start_year: int, start_month: int, leads: list[int]):
    if not leads:
        return []
    max_lead = max(leads)
    all_months = get_next_months(start_year, start_month, max_lead + 1)
    return [all_months[l] for l in leads]

def format_month_range(months: list[dict]) -> str:
    """Standard month-range style for ALL slides"""
    if not months:
        return ""
    start = months[0]
    end = months[-1]
    sep = " – "
    if start["buddhist_year"] == end["buddhist_year"]:
        return f"{start['thai_name']}{sep}{end['thai_name']} {start['buddhist_year']}"
    return (
        f"{start['thai_name']} {start['buddhist_year']}"
        f"{sep}"
        f"{end['thai_name']} {end['buddhist_year']}"
    )

def format_month_range_long(months: list[dict]) -> str:
    """Long-form range used on cover/summary: 'เดือนมีนาคม-สิงหาคม ปี 2569'"""
    if not months:
        return ""
    start = months[0]
    end = months[-1]
    if start["buddhist_year"] == end["buddhist_year"]:
        return f"เดือน{start['thai_name']}-{end['thai_name']} ปี {start['buddhist_year']}"
    return (
        f"เดือน{start['thai_name']} ปี {start['buddhist_year']}"
        f" – "
        f"{end['thai_name']} ปี {end['buddhist_year']}"
    )


# ==========================================
# 2. ฟังก์ชันจัดการ Text Box บน PowerPoint
# ==========================================
def _find_shape_recursive(shapes, shape_name: str):
    """ค้นหา Shape ตามชื่อแบบเจาะลึกเข้าไปใน Group"""
    for shape in shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found_shape = _find_shape_recursive(shape.shapes, shape_name)
            if found_shape:
                return found_shape
    return None

def replace_text_by_name(slide: Slide, shape_name: str, new_text: str) -> bool:
    """
    ค้นหาและแทนที่ข้อความตามชื่อ Object โดยรักษา Format (Font, Size, Color) เดิมไว้
    *แก้ปัญหา Runs พังโดยการเก็บรูปแบบของ Run แรกไว้*
    """
    target_shape = _find_shape_recursive(slide.shapes, shape_name)

    if not target_shape:
        logger.warning(f"ไม่พบ Text Box ชื่อ '{shape_name}'")
        return False

    if not target_shape.has_text_frame:
        logger.warning(f"Object '{shape_name}' ไม่ใช่กล่องข้อความ")
        return False

    logger.info(f"อัปเดตข้อความ '{shape_name}' -> '{new_text}'")
    
    text_frame = target_shape.text_frame

    # เทคนิครักษา Format: ใส่ข้อความใหม่ลงใน Run แรกสุด แล้วเคลียร์ Run อื่นๆ ทิ้งทั้งหมด
    if text_frame.paragraphs:
        p0 = text_frame.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = str(new_text)
            
            # ลบข้อความที่เหลือในบรรทัดเดียวกัน
            for run in p0.runs[1:]:
                run.text = ""
                
            # ลบข้อความในบรรทัดอื่นๆ (ถ้ามีหลายย่อหน้า)
            for paragraph in text_frame.paragraphs[1:]:
                for run in paragraph.runs:
                    run.text = ""
        else:
            # ถ้ามี Paragraph แต่ไม่มี Run เลย ให้ตั้งค่าตรงๆ
            target_shape.text = str(new_text)
    else:
        # ถ้าไม่มีอะไรเลย
        target_shape.text = str(new_text)

    return True


def replace_text_paragraphs(slide: Slide, shape_name: str, lines: list[str]) -> bool:
    """
    แทนที่ข้อความใน Text Box ด้วยหลาย Paragraph โดยรักษา Format ของ Paragraph แรกไว้
    ใช้สำหรับ txt_tbl_title ที่ต้องการ 2 บรรทัดจริงๆ (ไม่ใช่แค่ word-wrap)
    """
    target_shape = _find_shape_recursive(slide.shapes, shape_name)
    if not target_shape:
        logger.warning(f"ไม่พบ Text Box ชื่อ '{shape_name}'")
        return False
    if not target_shape.has_text_frame:
        logger.warning(f"Object '{shape_name}' ไม่ใช่กล่องข้อความ")
        return False

    tf = target_shape.text_frame
    if not tf.paragraphs:
        return False

    p0 = tf.paragraphs[0]

    # Set paragraph 0 to lines[0]
    line0 = lines[0] if lines else ""
    if p0.runs:
        p0.runs[0].text = line0
        for run in p0.runs[1:]:
            run.text = ""
    else:
        target_shape.text = line0

    # Add / update subsequent paragraphs
    for i, line in enumerate(lines[1:], start=1):
        existing = tf.paragraphs  # re-read after each insert
        if i < len(existing):
            pi = existing[i]
            if pi.runs:
                pi.runs[0].text = line
                for run in pi.runs[1:]:
                    run.text = ""
        else:
            # Deep-copy paragraph 0 XML, set new text, insert after last paragraph
            new_p = deepcopy(p0._p)
            r_els = new_p.findall(f"{{{_A_NS}}}r")
            if r_els:
                t_el = r_els[0].find(f"{{{_A_NS}}}t")
                if t_el is not None:
                    t_el.text = line
                for r in r_els[1:]:
                    t_sub = r.find(f"{{{_A_NS}}}t")
                    if t_sub is not None:
                        t_sub.text = ""
            tf.paragraphs[-1]._p.addnext(new_p)

    # Clear runs in any extra paragraphs beyond what we need
    for extra_p in tf.paragraphs[len(lines):]:
        for run in extra_p.runs:
            run.text = ""

    logger.info(f"อัปเดตข้อความ '{shape_name}' -> {lines}")
    return True


def update_date_runs(slide: Slide, shape_name: str, day: int, month_idx: int, year: int) -> bool:
    """
    Update the 3-run date line inside txt_report_date while preserving each run's
    language ID and font formatting.

    Expects paragraph index 2 to have exactly 3 runs:
      run 0: day number + space  (ENGLISH_US)
      run 1: full Thai month + space  (THAI)
      run 2: 4-digit Buddhist year  (ENGLISH_UK)
    """
    target_shape = _find_shape_recursive(slide.shapes, shape_name)
    if not target_shape:
        logger.warning(f"ไม่พบ shape '{shape_name}'")
        return False
    if not target_shape.has_text_frame:
        logger.warning(f"Shape '{shape_name}' ไม่ใช่กล่องข้อความ")
        return False

    paras = target_shape.text_frame.paragraphs
    if len(paras) < 3:
        logger.warning(f"'{shape_name}' มี {len(paras)} paragraphs (ต้องการอย่างน้อย 3)")
        return False

    p = paras[2]
    if len(p.runs) < 3:
        logger.warning(f"'{shape_name}' paragraph 2 มี {len(p.runs)} runs (ต้องการ 3)")
        return False

    p.runs[0].text = f"{day} "
    p.runs[1].text = f"{get_thai_month(month_idx)} "
    p.runs[2].text = str(get_buddhist_year(year))

    logger.info(f"อัปเดตวันที่ '{shape_name}' -> {day} {get_thai_month(month_idx)} {get_buddhist_year(year)}")
    return True