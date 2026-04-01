# src/manager.py
import copy
import logging
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.slide import Slide

logger = logging.getLogger(__name__)

# Relationship types used in slides
_LAYOUT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _patch_rids(element, rId_map: dict) -> None:
    """Walk the element tree and update all r:embed / r:id / r:link attribute values."""
    if not rId_map:
        return
    for attr in (f"{{{_R_NS}}}embed", f"{{{_R_NS}}}id", f"{{{_R_NS}}}link"):
        for el in element.iter():
            val = el.get(attr)
            if val in rId_map:
                el.set(attr, rId_map[val])


def _replace_element_inplace(target, source) -> None:
    """
    Copy source element's content into target, preserving target's object identity.
    Needed because slide.element and slide.part._element may be separate references —
    assigning new_part._element = new_element would leave slide.element stale.
    """
    target.clear()
    target.tag = source.tag
    target.text = source.text
    target.tail = source.tail
    for k, v in source.attrib.items():
        target.set(k, v)
    for child in source:
        target.append(copy.deepcopy(child))


class ReportManager:
    """
    ผู้จัดการไฟล์ PowerPoint
    รับผิดชอบการโหลด Template, ค้นหาสไลด์ (Tag), โคลน/ลบสไลด์, และบันทึกไฟล์
    """

    def __init__(self, template_path: Path | str):
        self.template_path = Path(template_path)
        self.prs = None

    # ─────────────────────────────────────────
    # Load / Save
    # ─────────────────────────────────────────

    def load_template(self) -> bool:
        """โหลดไฟล์ Template เข้าสู่หน่วยความจำ"""
        if not self.template_path.exists():
            logger.error(f"Template file not found: {self.template_path}")
            return False
        try:
            self.prs = Presentation(self.template_path)
            logger.info(f"Loaded template: {self.template_path.name}")
            return True
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            return False

    def save_report(self, final_output_path: Path | str) -> bool:
        """บันทึกไฟล์สไลด์ตาม Path ที่ได้มาจาก OutputManager"""
        if not self.prs:
            logger.error("Cannot save: Presentation is not loaded.")
            return False
        try:
            self.prs.save(str(final_output_path))
            logger.info(f"Saved report to: {final_output_path}")
            return True
        except Exception as e:
            logger.error(f"Failed to save report: {e}")
            return False

    # ─────────────────────────────────────────
    # Slide lookup
    # ─────────────────────────────────────────

    def get_slide_by_tag(self, tag_name: str) -> Slide | None:
        """ค้นหาสไลด์ที่มี Object ชื่อตรงกับ tag_name"""
        if not self.prs:
            logger.error("Cannot search for tags: Presentation is not loaded.")
            return None
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name == tag_name:
                    logger.debug(f"Found tag '{tag_name}' on slide {self._slide_index(slide) + 1}")
                    return slide
        logger.warning(f"Slide with tag '{tag_name}' not found in the template.")
        return None

    # ─────────────────────────────────────────
    # Slide cloning
    # ─────────────────────────────────────────

    def clone_slide_after(self, source_slide: Slide, after_slide: Slide) -> Slide:
        """
        Clone source_slide (copying all relationships and XML) and insert the
        new slide immediately after after_slide in the presentation.
        """
        source_part = source_slide.part

        # 1. Add a fresh slide at the end using the same layout
        new_slide = self.prs.slides.add_slide(source_slide.slide_layout)
        new_part = new_slide.part

        # 2. Map source layout rId → new slide's layout rId
        rId_map: dict[str, str] = {}
        source_layout_rId = next(
            (rId for rId, rel in source_part.rels.items()
             if rel.reltype == _LAYOUT_RELTYPE), None
        )
        new_layout_rId = next(
            (rId for rId, rel in new_part.rels.items()
             if rel.reltype == _LAYOUT_RELTYPE), None
        )
        if source_layout_rId and new_layout_rId:
            rId_map[source_layout_rId] = new_layout_rId

        # 3. Copy all other (non-layout) internal relationships
        for rId, rel in source_part.rels.items():
            if rel.is_external or rel.reltype == _LAYOUT_RELTYPE:
                continue
            new_rId = new_part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

        # 4. Deep-copy source XML, patch rId attributes, assign to both
        # part and slide proxy, then clear the lazyproperty cache so
        # new_slide.shapes reads from the new element (not the cached empty one).
        new_element = copy.deepcopy(source_part._element)
        _patch_rids(new_element, rId_map)
        new_part._element = new_element
        new_slide._element = new_element
        new_slide.__dict__.pop("shapes", None)

        # 5. Move the new slide to immediately after after_slide
        target_idx = self._slide_index(after_slide) + 1
        self._move_slide_to_idx(new_slide, target_idx)

        logger.debug(f"Cloned slide inserted at index {target_idx}")
        return new_slide

    # ─────────────────────────────────────────
    # Slide removal
    # ─────────────────────────────────────────

    def remove_slide(self, slide: Slide) -> None:
        """Remove a slide from the presentation.

        Removes the slide from _sldIdLst AND drops the presentation's
        relationship to it, so the part is unreachable from iter_parts()
        and won't be written to the zip (preventing duplicate-partname warnings).
        """
        idx = self._slide_index(slide)
        if idx < 0:
            logger.warning("Cannot remove slide: not found in presentation.")
            return

        slide_part = slide.part
        prs_part   = self.prs.part

        # Find the rId in prs.part that points to this slide
        rId_to_remove = None
        for rId, rel in prs_part.rels.items():
            if not rel.is_external and rel.target_part is slide_part:
                rId_to_remove = rId
                break

        # Remove from _sldIdLst XML
        xml_slides = self.prs.slides._sldIdLst
        xml_slides.remove(xml_slides[idx])

        # Drop the relationship so the part is no longer reachable
        if rId_to_remove:
            prs_part.drop_rel(rId_to_remove)

        logger.info(f"Removed slide at index {idx}")

    # ─────────────────────────────────────────
    # Internal helpers
    # ─────────────────────────────────────────

    def _slide_index(self, slide: Slide) -> int:
        """Return 0-based index of slide in the presentation (-1 if not found)."""
        try:
            return list(self.prs.slides).index(slide)
        except ValueError:
            return -1

    def _move_slide_to_idx(self, slide: Slide, target_idx: int) -> None:
        """Move slide to the given 0-based position in the presentation."""
        current_idx = self._slide_index(slide)
        if current_idx == target_idx:
            return
        xml_slides = self.prs.slides._sldIdLst
        sld_id = xml_slides[current_idx]
        xml_slides.remove(sld_id)
        xml_slides.insert(target_idx, sld_id)
